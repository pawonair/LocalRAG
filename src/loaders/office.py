"""
Office Document Loaders Module
Handles Word, Excel, and PowerPoint files.
"""

import tempfile
from pathlib import Path
from typing import List

from .base import BaseLoader, LoaderResult, Document


class WordLoader(BaseLoader):
    """Loader for Microsoft Word documents (.docx, .doc)."""

    SUPPORTED_EXTENSIONS = [".docx", ".doc"]

    def load(self, file_path: str) -> LoaderResult:
        """Load a Word document."""
        try:
            from docx import Document as DocxDocument

            doc = DocxDocument(file_path)
            documents = []

            # Extract paragraphs
            current_section = []
            current_heading = ""

            for para in doc.paragraphs:
                text = para.text.strip()

                if not text:
                    continue

                # Check if it's a heading
                if para.style.name.startswith("Heading"):
                    # Save previous section
                    if current_section:
                        content = "\n".join(current_section)
                        documents.append(Document(
                            page_content=content,
                            metadata={
                                **self._create_metadata(file_path),
                                "section": current_heading or "Introduction",
                            }
                        ))
                        current_section = []

                    current_heading = text
                    current_section.append(f"## {text}")
                else:
                    current_section.append(text)

            # Add last section
            if current_section:
                content = "\n".join(current_section)
                documents.append(Document(
                    page_content=content,
                    metadata={
                        **self._create_metadata(file_path),
                        "section": current_heading or "Content",
                    }
                ))

            # If no documents created, extract all text
            if not documents:
                full_text = "\n".join(p.text for p in doc.paragraphs if p.text.strip())
                documents.append(Document(
                    page_content=full_text,
                    metadata=self._create_metadata(file_path)
                ))

            # Also extract tables
            for i, table in enumerate(doc.tables):
                table_text = self._extract_table(table)
                if table_text:
                    documents.append(Document(
                        page_content=table_text,
                        metadata={
                            **self._create_metadata(file_path),
                            "type": "table",
                            "table_index": i + 1,
                        }
                    ))

            return LoaderResult(
                success=True,
                documents=documents,
                file_type="word",
                metadata={"paragraph_count": len(doc.paragraphs), "table_count": len(doc.tables)}
            )

        except ImportError:
            return LoaderResult(
                success=False,
                error="python-docx is required for Word documents. Install with: pip install python-docx",
                file_type="word"
            )
        except Exception as e:
            return LoaderResult(
                success=False,
                error=f"Failed to load Word document: {str(e)}",
                file_type="word"
            )

    def load_from_bytes(self, content: bytes, filename: str) -> LoaderResult:
        """Load Word document from bytes."""
        try:
            with tempfile.NamedTemporaryFile(suffix=".docx", delete=False) as tmp:
                tmp.write(content)
                tmp_path = tmp.name

            result = self.load(tmp_path)

            # Update metadata
            for doc in result.documents:
                doc.metadata["source"] = filename
                doc.metadata["filename"] = filename

            Path(tmp_path).unlink(missing_ok=True)
            return result

        except Exception as e:
            return LoaderResult(
                success=False,
                error=f"Failed to load Word document from bytes: {str(e)}",
                file_type="word"
            )

    def _extract_table(self, table) -> str:
        """Extract table content as formatted text."""
        rows = []
        for row in table.rows:
            cells = [cell.text.strip() for cell in row.cells]
            rows.append(" | ".join(cells))

        if rows:
            # Add header separator
            if len(rows) > 1:
                header = rows[0]
                separator = " | ".join(["---"] * len(rows[0].split(" | ")))
                return "\n".join([header, separator] + rows[1:])

        return "\n".join(rows)


class ExcelLoader(BaseLoader):
    """Loader for Microsoft Excel files (.xlsx, .xls)."""

    SUPPORTED_EXTENSIONS = [".xlsx", ".xls"]

    def load(self, file_path: str) -> LoaderResult:
        """Load an Excel file."""
        try:
            import openpyxl

            wb = openpyxl.load_workbook(file_path, data_only=True)
            documents = []

            for sheet_name in wb.sheetnames:
                sheet = wb[sheet_name]

                # Extract data from sheet
                rows = []
                headers = []

                for i, row in enumerate(sheet.iter_rows(values_only=True)):
                    # Filter out empty rows
                    if any(cell is not None for cell in row):
                        if i == 0:
                            headers = [str(c) if c else f"Column_{j}" for j, c in enumerate(row)]
                        else:
                            rows.append(row)

                if not rows:
                    continue

                # Format as text
                ROWS_PER_DOC = 50

                for chunk_start in range(0, len(rows), ROWS_PER_DOC):
                    chunk_rows = rows[chunk_start:chunk_start + ROWS_PER_DOC]

                    lines = [f"Sheet: {sheet_name}", f"Headers: {', '.join(headers)}", ""]

                    for row in chunk_rows:
                        row_data = []
                        for j, cell in enumerate(row):
                            header = headers[j] if j < len(headers) else f"Column_{j}"
                            row_data.append(f"{header}: {cell}")
                        lines.append(" | ".join(row_data))

                    documents.append(Document(
                        page_content="\n".join(lines),
                        metadata={
                            **self._create_metadata(file_path),
                            "sheet": sheet_name,
                            "headers": headers,
                            "row_start": chunk_start,
                            "row_end": min(chunk_start + ROWS_PER_DOC, len(rows)),
                        }
                    ))

            wb.close()

            if not documents:
                return LoaderResult(
                    success=False,
                    error="No data found in Excel file",
                    file_type="excel"
                )

            return LoaderResult(
                success=True,
                documents=documents,
                file_type="excel",
                metadata={"sheet_count": len(wb.sheetnames)}
            )

        except ImportError:
            return LoaderResult(
                success=False,
                error="openpyxl is required for Excel files. Install with: pip install openpyxl",
                file_type="excel"
            )
        except Exception as e:
            return LoaderResult(
                success=False,
                error=f"Failed to load Excel file: {str(e)}",
                file_type="excel"
            )

    def load_from_bytes(self, content: bytes, filename: str) -> LoaderResult:
        """Load Excel file from bytes."""
        try:
            with tempfile.NamedTemporaryFile(suffix=".xlsx", delete=False) as tmp:
                tmp.write(content)
                tmp_path = tmp.name

            result = self.load(tmp_path)

            for doc in result.documents:
                doc.metadata["source"] = filename
                doc.metadata["filename"] = filename

            Path(tmp_path).unlink(missing_ok=True)
            return result

        except Exception as e:
            return LoaderResult(
                success=False,
                error=f"Failed to load Excel from bytes: {str(e)}",
                file_type="excel"
            )


class PowerPointLoader(BaseLoader):
    """Loader for Microsoft PowerPoint files (.pptx, .ppt)."""

    SUPPORTED_EXTENSIONS = [".pptx", ".ppt"]

    def load(self, file_path: str) -> LoaderResult:
        """Load a PowerPoint file."""
        try:
            from pptx import Presentation

            prs = Presentation(file_path)
            documents = []

            for i, slide in enumerate(prs.slides):
                slide_content = []
                slide_content.append(f"=== Slide {i + 1} ===")

                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_content.append(shape.text.strip())

                    # Extract table content
                    if shape.has_table:
                        table_text = self._extract_table(shape.table)
                        if table_text:
                            slide_content.append(f"\n[Table]\n{table_text}")

                if len(slide_content) > 1:  # More than just the slide header
                    documents.append(Document(
                        page_content="\n\n".join(slide_content),
                        metadata={
                            **self._create_metadata(file_path),
                            "slide_number": i + 1,
                            "total_slides": len(prs.slides),
                        }
                    ))

            if not documents:
                return LoaderResult(
                    success=False,
                    error="No text content found in PowerPoint file",
                    file_type="powerpoint"
                )

            return LoaderResult(
                success=True,
                documents=documents,
                file_type="powerpoint",
                metadata={"slide_count": len(prs.slides)}
            )

        except ImportError:
            return LoaderResult(
                success=False,
                error="python-pptx is required for PowerPoint files. Install with: pip install python-pptx",
                file_type="powerpoint"
            )
        except Exception as e:
            return LoaderResult(
                success=False,
                error=f"Failed to load PowerPoint file: {str(e)}",
                file_type="powerpoint"
            )

    def load_from_bytes(self, content: bytes, filename: str) -> LoaderResult:
        """Load PowerPoint from bytes."""
        try:
            with tempfile.NamedTemporaryFile(suffix=".pptx", delete=False) as tmp:
                tmp.write(content)
                tmp_path = tmp.name

            result = self.load(tmp_path)

            for doc in result.documents:
                doc.metadata["source"] = filename
                doc.metadata["filename"] = filename

            Path(tmp_path).unlink(missing_ok=True)
            return result

        except Exception as e:
            return LoaderResult(
                success=False,
                error=f"Failed to load PowerPoint from bytes: {str(e)}",
                file_type="powerpoint"
            )

    def _extract_table(self, table) -> str:
        """Extract table content as formatted text."""
        rows = []
        for row in table.rows:
            cells = [cell.text.strip() for cell in row.cells]
            rows.append(" | ".join(cells))

        return "\n".join(rows)
